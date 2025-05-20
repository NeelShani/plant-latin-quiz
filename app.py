import streamlit as st
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
import io, random

# ─── Config & Welcome ───────────────────────────────────────────
st.set_page_config(page_title="Plant Name Memory Quiz", layout="wide")
st.title("Plant Quiz")
st.markdown("### Hello Anna! Welcome to this specialized app created just for you.")

# ─── Constants ──────────────────────────────────────────────────
FIXED_HEIGHT = 600  # px for every image

# ─── State Defaults ─────────────────────────────────────────────
for key, default in {
    "plants":     [],
    "subset":     [],
    "remaining":  [],
    "current":    None,
    "revealed":   False,
    "guess_input":""
}.items():
    if key not in st.session_state:
        st.session_state[key] = default

# ─── Callbacks ──────────────────────────────────────────────────
def start_quiz():
    # slice out chosen range (or all)
    total = len(st.session_state.plants)
    if st.session_state.quiz_all:
        s, e = 1, total
    else:
        s, e = st.session_state.range_start, st.session_state.range_end

    sub = st.session_state.plants[s-1:e]
    st.session_state.subset = sub
    st.session_state.remaining = list(range(len(sub)))
    random.shuffle(st.session_state.remaining)
    st.session_state.current = None
    st.session_state.revealed = False
    st.session_state.guess_input = ""

def next_slide():
    st.session_state.current = None
    st.session_state.revealed = False
    st.session_state.guess_input = ""

def restart_quiz():
    st.session_state.remaining = list(range(len(st.session_state.subset)))
    random.shuffle(st.session_state.remaining)
    st.session_state.current = None
    st.session_state.revealed = False
    st.session_state.guess_input = ""

def maybe_reveal():
    """If user has typed a non-empty guess and hits Enter, reveal answer."""
    if st.session_state.guess_input.strip():
        st.session_state.revealed = True

# ─── Extraction ─────────────────────────────────────────────────
def extract_slides(pptx_bytes):
    prs = Presentation(io.BytesIO(pptx_bytes))
    slides = []
    for slide in prs.slides:
        # find first picture
        img = None
        for shp in slide.shapes:
            if shp.shape_type in (MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.PLACEHOLDER):
                try:
                    blob = shp.image.blob
                    img = Image.open(io.BytesIO(blob))
                    break
                except Exception:
                    pass
        # gather all text
        lines = [
            shp.text.strip() for shp in slide.shapes
            if hasattr(shp, "text") and shp.text.strip()
        ]
        if img and lines:
            slides.append((img, "\n".join(lines)))
    return slides

# ─── UI: File Upload & Quiz Setup ───────────────────────────────
pptx_file = st.file_uploader("Upload your PowerPoint (.pptx)", type=["pptx"])
if pptx_file:
    st.session_state.plants = extract_slides(pptx_file.read())
    total = len(st.session_state.plants)
    if total == 0:
        st.error("❌ No valid slides found. Make sure each slide has one image + ≥1 text.")
        st.stop()
    st.success(f"✓ Loaded {total} slides.")

    st.session_state.quiz_all = st.checkbox("Quiz **all** slides (ignore range)", value=True)
    if not st.session_state.quiz_all:
        start, end = st.slider("Select slide‐range to quiz:", 1, total, (1, total))
        st.session_state.range_start = start
        st.session_state.range_end   = end

    st.button("Start Quiz", on_click=start_quiz)

# ─── UI: Quiz Loop ────────────────────────────────────────────────
if st.session_state.remaining:
    # pick new slide if needed
    if st.session_state.current is None:
        st.session_state.current = st.session_state.remaining.pop()
    img, text = st.session_state.subset[st.session_state.current]

    # resize to fixed height
    w, h = img.size
    new_w = int(w * FIXED_HEIGHT / h)
    resized = img.resize((new_w, FIXED_HEIGHT), Image.LANCZOS)

    st.image(resized, use_column_width=False)

    # guess box with on_change callback
    st.text_input(
        "Your guess (optional):",
        key="guess_input",
        on_change=maybe_reveal
    )

    # buttons: reveal only if no guess yet, next always
    cols = st.columns(2)
    with cols[0]:
        if (not st.session_state.revealed) and (not st.session_state.guess_input.strip()):
            st.button("Reveal Answer", on_click=lambda: st.session_state.update({"revealed": True}))
    with cols[1]:
        st.button("Next", on_click=next_slide)

    # show answer once revealed or guess_entered
    if st.session_state.revealed:
        st.markdown("**Answer (all text on slide):**")
        st.write(text)

elif st.session_state.subset:
    st.info("✅ You've seen **all** slides in this range!")
    st.button("Restart Quiz", on_click=restart_quiz)

st.markdown("### Good luck for the exam — you can do it! ❤️")
